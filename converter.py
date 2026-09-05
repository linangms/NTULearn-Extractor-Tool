import io
import os
import re
import tempfile
import zipfile
import logging
from typing import Any, Callable, Dict, List, Optional, Tuple
from bs4 import BeautifulSoup
import markdownify

from blackboard_client import FileTooLargeError, DiskFile, _peek, _content_len, _discard

logger = logging.getLogger("converter")


def _format_mb(size_bytes: int) -> str:
    return f"{size_bytes / (1024 * 1024):.0f}MB"

def _is_valid_video_bytes(data) -> bool:
    """
    Checks for a real video binary signature so an HTML error/login page
    (e.g. from an unresolved LTI launch link) doesn't get saved as if it
    were the video itself. A DiskFile (content streamed straight to disk
    because it grew past the in-memory spool threshold) is always genuine
    binary media by construction - an HTML error page never gets that
    large - so it's treated as valid without re-reading it from disk.
    """
    if isinstance(data, DiskFile):
        return True
    if not data or len(data) < 1000:
        return False
    snippet = data[:64]
    if snippet.startswith(b"<") or snippet.startswith(b"{") or b"<?xml" in snippet or b"404 Not Found" in snippet:
        return False
    return bool(b"ftyp" in snippet or b"moov" in snippet or snippet.startswith(b"\x00\x00\x00") or b"matroska" in snippet)


def _write_downloaded(zf: "zipfile.ZipFile", zip_target_path: str, data) -> Optional[bytes]:
    """
    Writes a downloaded payload (raw bytes, or a DiskFile streamed straight
    to disk to avoid buffering a whole video in memory) into the zip.
    Returns the bytes for small in-memory content, so callers can still
    sniff or parse it (captions, text extraction); returns None for
    on-disk content, which is always large binary media - not something
    this tool parses text from or needs to caption-sniff.
    """
    if isinstance(data, DiskFile):
        zf.write(data.path, zip_target_path)
        _discard(data)
        return None
    zf.writestr(zip_target_path, data)
    return data


async def _ensure_expanded(item: Dict[str, Any], course_id: Optional[str], expand_fn: Optional[Callable]) -> Dict[str, Any]:
    """
    Expands a raw, unbuilt child item (as returned by Blackboard's shallow
    /children listing) into a fully built node - title, body, attachments,
    Kaltura resolution, and its OWN direct children as raw items in turn -
    by calling expand_fn (BlackboardClient.build_content_node) on demand,
    one node at a time, instead of a whole folder tree needing to already be
    built before any of it can be processed. Already-built nodes (recognized
    by an "attachments" key, which every real build sets and a raw API item
    never has) and pre-built trees (expand_fn is None, e.g. the mock/test
    content trees) pass through unchanged.
    """
    if expand_fn is not None and "attachments" not in item:
        return await expand_fn(course_id, item)
    return item


def _looks_like_kaltura_attachment(attachment_id: str, att: Dict[str, Any]) -> bool:
    """
    Cheap local check so we only ask the caption_downloader to hit Kaltura's API
    for attachments that are actually Kaltura video entries.
    """
    if att.get("isKaltura"):
        return True
    text = attachment_id or ""
    return "kaltura" in text.lower() or bool(re.search(r'\b[01]_[a-zA-Z0-9]{8,12}\b', text))


def sanitize_filename(name: str) -> str:
    """
    Sanitizes string for cross-platform safe file and folder names.
    """
    if not name:
        return "Untitled"
    # Replace illegal characters with underscore
    sanitized = re.sub(r'[\\/*?:"<>|]', "_", name)
    sanitized = re.sub(r"\s+", " ", sanitized).strip()
    return sanitized or "Untitled"


def clean_vtt_srt_transcript(text: str) -> str:
    """
    Cleans WebVTT (.vtt) and SubRip (.srt) subtitle content into readable plain text transcript.
    """
    if not text:
        return ""
    # Remove WEBVTT header
    text = re.sub(r'^WEBVTT.*?\n', '', text, flags=re.IGNORECASE)
    # Remove timestamp lines (e.g. 00:00:01.000 --> 00:00:04.000 or 00:00:01,000 --> 00:00:04,000)
    text = re.sub(r'\d{2}:\d{2}:\d{2}[\.,]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[\.,]\d{3}.*?\n', '', text)
    text = re.sub(r'\d{2}:\d{2}[\.,]\d{3}\s*-->\s*\d{2}:\d{2}[\.,]\d{3}.*?\n', '', text)
    # Remove standalone cue numbers
    text = re.sub(r'^\d+\s*$', '', text, flags=re.MULTILINE)
    # Remove HTML/XML cue tags like <v Speaker> or <c>
    text = re.sub(r'<[^>]+>', '', text)
    
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    dedup_lines = []
    for line in lines:
        if not dedup_lines or line != dedup_lines[-1]:
            dedup_lines.append(line)
    return "\n\n".join(dedup_lines)


def extract_text_from_file_bytes(file_bytes: bytes, filename: str) -> str:
    """
    Extracts readable text content from PDF, PPTX, DOCX, TXT, VTT, SRT, or HTML attachment files.
    """
    if not file_bytes:
        return ""

    lower_name = filename.lower()

    # 1. PDF extraction
    if lower_name.endswith(".pdf"):
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            pages_text = []
            for idx, page in enumerate(reader.pages, 1):
                t = page.extract_text() or ""
                if t.strip():
                    pages_text.append(f"#### Page {idx}\n\n{t.strip()}")
            if pages_text:
                return "\n\n".join(pages_text)
        except Exception as e:
            logger.warning(f"Could not extract text from PDF {filename}: {e}")

    # 2. PowerPoint (.pptx) extraction
    elif lower_name.endswith(".pptx"):
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for idx, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_content.append(shape.text.strip())
                if slide_content:
                    slides_text.append(f"#### Slide {idx}\n\n" + "\n\n".join(slide_content))
            if slides_text:
                return "\n\n".join(slides_text)
        except Exception as e:
            logger.warning(f"Could not extract text from PPTX {filename}: {e}")

    # 3. Word (.docx) extraction
    elif lower_name.endswith(".docx"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            if paras:
                return "\n\n".join(paras)
        except Exception as e:
            logger.warning(f"Could not extract text from DOCX {filename}: {e}")

    # 4. WebVTT / SRT Subtitle Transcript extraction
    elif any(lower_name.endswith(ext) for ext in [".vtt", ".srt"]):
        try:
            raw_text = file_bytes.decode("utf-8", errors="ignore")
            return clean_vtt_srt_transcript(raw_text)
        except Exception as e:
            logger.warning(f"Could not parse VTT/SRT transcript {filename}: {e}")

    # 5. Plain Text / HTML / Markdown / XML
    elif any(lower_name.endswith(ext) for ext in [".txt", ".html", ".htm", ".md", ".csv", ".json", ".xml"]):
        try:
            raw_text = file_bytes.decode("utf-8", errors="ignore")
            if lower_name.endswith(".html") or lower_name.endswith(".htm"):
                return markdownify.markdownify(raw_text)
            return raw_text
        except Exception as e:
            logger.warning(f"Could not decode text file {filename}: {e}")

    return ""


class CourseMarkdownConverter:
    """
    Parses Blackboard content trees, converts HTML into Markdown,
    downloads and places attachments into relative asset directories,
    rewrites links, and builds a downloadable ZIP file structure.
    """

    def __init__(self, course_name: str, course_id: str, base_url: str = "https://ntulearntst.ntu.edu.sg"):
        self.course_name = course_name or course_id or "Course Materials"
        self.course_id = course_id
        self.base_url = base_url.rstrip("/")
        clean_name = re.sub(r'\s*-\s*Course Materials$', '', self.course_name, flags=re.IGNORECASE).strip()
        self.clean_course_name = clean_name or self.course_name
        self.root_folder_name = sanitize_filename(f"{self.clean_course_name}_Markdown")

    def convert_html_to_markdown(self, html_content: str, attachment_map: Optional[Dict[str, str]] = None) -> str:
        if not html_content or not html_content.strip():
            return ""

        soup = BeautifulSoup(html_content, "html.parser")

        if attachment_map:
            for a_tag in soup.find_all("a", href=True):
                href = a_tag["href"]
                for orig_url, rel_path in attachment_map.items():
                    if orig_url in href or href in orig_url:
                        a_tag["href"] = rel_path

            for img_tag in soup.find_all("img", src=True):
                src = img_tag["src"]
                for orig_url, rel_path in attachment_map.items():
                    if orig_url in src or src in orig_url:
                        img_tag["src"] = rel_path

        md = markdownify.markdownify(str(soup), heading_style="ATX")
        md = re.sub(r"\n{3,}", "\n\n", md).strip()
        return md

    async def _flatten_tree(
        self,
        nodes: List[Dict[str, Any]],
        flat_list: List[Dict[str, Any]],
        folder_path: str = "",
        course_id: Optional[str] = None,
        expand_fn: Optional[Callable] = None,
    ):
        # expand_fn (BlackboardClient.build_content_node), when given, expands
        # a raw child stub into a full node one at a time, right before it's
        # needed - see _ensure_expanded. Without it (a pre-built tree, e.g.
        # the mock/test content trees), every node is already fully built and
        # this behaves exactly as it always did.
        for raw_node in nodes:
            node = await _ensure_expanded(raw_node, course_id, expand_fn)
            title = node.get("title", "Untitled")
            current_path = f"{folder_path} / {title}" if folder_path else title

            body = (
                node.get("body")
                or node.get("description")
                or node.get("instructions")
                or node.get("summary")
                or node.get("formattedBody")
                or ""
            )
            has_content = bool(body and body.strip()) or bool(node.get("attachments"))

            children = node.get("children", [])
            if children:
                if has_content:
                    node_copy = dict(node)
                    node_copy["folder_path"] = folder_path
                    flat_list.append(node_copy)
                await self._flatten_tree(children, flat_list, folder_path=current_path, course_id=course_id, expand_fn=expand_fn)
            else:
                node_copy = dict(node)
                node_copy["folder_path"] = folder_path
                flat_list.append(node_copy)

    async def _iter_flatten_tree(
        self,
        nodes: List[Dict[str, Any]],
        folder_path: str = "",
        course_id: Optional[str] = None,
        expand_fn: Optional[Callable] = None,
    ):
        """
        Same traversal as _flatten_tree, but yields one flattened item at a
        time instead of collecting the whole (possibly lazily-expanded)
        subtree into a list first. _flatten_tree's expand_fn support alone
        isn't enough to keep a big folder tree's memory bounded, since it
        still builds a complete flat_list before returning - this generator
        is what lets build_zip_package_streaming release each item after
        it's written, one at a time, all the way down into nested folders,
        not just between top-level topics.
        """
        for raw_node in nodes:
            node = await _ensure_expanded(raw_node, course_id, expand_fn)
            title = node.get("title", "Untitled")
            current_path = f"{folder_path} / {title}" if folder_path else title

            body = (
                node.get("body")
                or node.get("description")
                or node.get("instructions")
                or node.get("summary")
                or node.get("formattedBody")
                or ""
            )
            has_content = bool(body and body.strip()) or bool(node.get("attachments"))

            children = node.get("children", [])
            if children:
                if has_content:
                    node_copy = dict(node)
                    node_copy["folder_path"] = folder_path
                    yield node_copy
                async for item in self._iter_flatten_tree(children, folder_path=current_path, course_id=course_id, expand_fn=expand_fn):
                    yield item
            else:
                node_copy = dict(node)
                node_copy["folder_path"] = folder_path
                yield node_copy

    def _build_readme(self, index_entries: List[Tuple[str, str, str]]) -> str:
        """index_entries: list of (title, folder_path, filename) - cheap metadata only."""
        readme_content = f"# {self.course_name}\n\n"
        readme_content += f"**Course ID:** {self.course_id}\n\n"
        readme_content += "Extracted & converted to Markdown via **NTULearn Extractor Tool**.\n\n"
        readme_content += "## Course Content Index\n\n"
        for title, folder_path, filename in index_entries:
            mod_path = f"*(Module: `{folder_path}`)* " if folder_path else ""
            readme_content += f"- [{title}]({filename}) {mod_path}\n"
        return readme_content

    async def _write_markdown_item(
        self,
        item: Dict[str, Any],
        idx: int,
        out_folder: str,
        zf: zipfile.ZipFile,
        attachment_downloader: Optional[Callable],
        caption_downloader: Optional[Callable],
        progress_callback: Optional[Callable],
        pct: float,
    ) -> Tuple[str, str, str]:
        """
        Converts one content-tree item to a Markdown file inside the
        already-open zip, handling its attachments/captions/video-transcript
        fallback. Returns (filename, title, folder_path) - the lightweight
        tuple needed for the course-level README index, so a caller doesn't
        need to hold onto the item's heavy body text after this returns.
        """
        clean_title = sanitize_filename(item.get("title", f"Item_{idx}"))
        filename = f"{idx:02d}_{clean_title}.md"
        file_path = f"{out_folder}/{filename}"

        if progress_callback:
            await progress_callback(f"Converting Markdown for: {item.get('title')}", pct)

        attachments_dir = f"{out_folder}/attachments"
        attachment_map, doc_text_map = await self._handle_attachments(
            item, attachments_dir, zf, attachment_downloader,
            caption_downloader=caption_downloader, progress_callback=progress_callback, pct=pct,
        )

        body = (
            item.get("body")
            or item.get("description")
            or item.get("instructions")
            or item.get("summary")
            or item.get("formattedBody")
            or ""
        )

        handler = item.get("contentHandler", {})
        handler_id = handler.get("id", "") if isinstance(handler, dict) else str(handler)
        is_video_item = (
            "video" in handler_id.lower()
            or "kaltura" in handler_id.lower()
            or "panopto" in handler_id.lower()
            or "media" in handler_id.lower()
            or any(ext in body.lower() for ext in [".mp4", ".mov", ".m4v", ".webm"])
        )

        # If this is a video item, try to fetch its real Kaltura captions first;
        # only fall back to the description text as a transcript if none exist.
        if is_video_item:
            clean_item_title = sanitize_filename(item.get("title", "Video"))
            txt_filename = f"{clean_item_title}_transcript.txt"
            if txt_filename not in doc_text_map:
                real_transcript_saved = False
                if caption_downloader:
                    search_text = body + " " + " ".join(
                        f"{att.get('originalUrl') or ''} {att.get('downloadUrl') or ''}"
                        for att in item.get("attachments", [])
                    )
                    if _looks_like_kaltura_attachment(search_text, {}):
                        try:
                            cap_data = await caption_downloader(self.course_id, item.get("id"), search_text)
                            if cap_data:
                                self._save_caption_and_transcript(item, cap_data, "captions.srt", zf, out_folder, doc_text_map)
                                real_transcript_saved = True
                        except Exception as e:
                            logger.warning(f"Failed to fetch Kaltura captions for {clean_item_title}: {e}")

                if not real_transcript_saved:
                    clean_body_text = BeautifulSoup(body, "html.parser").get_text(separator="\n\n", strip=True) if body else ""
                    if clean_body_text and len(clean_body_text) > 10:
                        zf.writestr(f"{out_folder}/{txt_filename}", clean_body_text)
                        doc_text_map[txt_filename] = clean_body_text

        md_text = self.convert_html_to_markdown(body, attachment_map)

        full_md = f"# {item.get('title')}\n\n"
        if item.get("folder_path"):
            full_md += f"**Module / Path:** `{item.get('folder_path')}`\n\n"

        if md_text and md_text.strip():
            full_md += f"{md_text.strip()}\n\n"

        # Append text extracted directly from attached PDF / PPTX / DOCX / TXT / Video Transcripts!
        if doc_text_map:
            for att_file, doc_text in doc_text_map.items():
                section_heading = "Video Transcript" if "transcript" in att_file.lower() or att_file.endswith(".txt") else "Extracted Content from Document"
                full_md += f"## {section_heading} (`{att_file}`)\n\n"
                full_md += f"{doc_text}\n\n"

        if item.get("attachments"):
            full_md += "### Attached Files & Resources\n\n"
            for att in item["attachments"]:
                att_name = sanitize_filename(att.get("fileName", "attachment"))
                rel_link = f"./attachments/{att_name}"
                full_md += f"- 📎 [{att.get('fileName', 'attachment')}]({rel_link})\n"

        zf.writestr(file_path, full_md)

        return filename, item.get("title"), item.get("folder_path")

    async def build_zip_package(
        self,
        content_tree: List[Dict[str, Any]],
        attachment_downloader: Optional[Callable[[str, str, str], Any]] = None,
        caption_downloader: Optional[Callable[[str, str, str], Any]] = None,
        progress_callback: Optional[Callable[[str, float], Any]] = None,
    ) -> str:
        """
        Processes content tree, creates all converted .md files in ONE single folder inside the zip archive.
        Writes directly to a temp file on disk (rather than an in-memory buffer) so a
        large course doesn't hold its whole compressed archive in RAM at once. Returns
        the temp file path; the caller is responsible for deleting it once served.
        """
        fd, zip_path = tempfile.mkstemp(suffix=".zip")
        os.close(fd)

        flat_items = []
        await self._flatten_tree(content_tree, flat_items)
        total_nodes = len(flat_items)

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            out_folder = self.root_folder_name
            index_entries = []

            for idx, item in enumerate(flat_items, 1):
                pct = (idx / max(1, total_nodes)) * 100
                filename, title, folder_path = await self._write_markdown_item(
                    item, idx, out_folder, zf, attachment_downloader, caption_downloader, progress_callback, pct
                )
                index_entries.append((title, folder_path, filename))

            zf.writestr(f"{out_folder}/00_README.md", self._build_readme(index_entries))

        return zip_path

    async def build_zip_package_streaming(
        self,
        topics,
        total_topics: int,
        attachment_downloader: Optional[Callable[[str, str, str], Any]] = None,
        caption_downloader: Optional[Callable[[str, str, str], Any]] = None,
        progress_callback: Optional[Callable[[str, float], Any]] = None,
        course_id: Optional[str] = None,
        expand_fn: Optional[Callable] = None,
    ) -> str:
        """
        Same output as build_zip_package, but `topics` is an async iterator
        that yields one top-level topic node at a time instead of a
        pre-built list of the whole course, and (when expand_fn is given)
        each topic's own folders are expanded and released one at a time too
        via _iter_flatten_tree, instead of requiring a topic's entire nested
        subtree to be fetched before any of it can be processed. This is
        what keeps a single unusually large topic (e.g. one folder holding
        many tutorials' worth of PDFs and videos) from being just as much of
        a memory spike as the whole course used to be.
        """
        fd, zip_path = tempfile.mkstemp(suffix=".zip")
        os.close(fd)
        out_folder = self.root_folder_name
        index_entries = []
        idx = 0

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            topic_num = 0
            async for top_node in topics:
                topic_num += 1
                # Rough, cheap estimate of this topic's size (top_node plus
                # whatever direct children are already visible) purely to
                # interpolate progress within the topic - it undercounts
                # deeply nested folders, but doesn't require expanding
                # anything to compute, so it doesn't cost the memory it's
                # trying to help report on.
                topic_item_count = self._count_nodes([top_node])
                i = 0
                async for item in self._iter_flatten_tree([top_node], course_id=course_id, expand_fn=expand_fn):
                    idx += 1
                    i += 1
                    within_topic = min(i / max(1, topic_item_count), 1.0)
                    pct = ((topic_num - 1 + within_topic) / max(1, total_topics)) * 100
                    filename, title, folder_path = await self._write_markdown_item(
                        item, idx, out_folder, zf, attachment_downloader, caption_downloader, progress_callback, pct
                    )
                    index_entries.append((title, folder_path, filename))
                # top_node and everything _iter_flatten_tree expanded from it
                # fall out of scope here, freeing this topic's body text
                # before the next topic is fetched.

            zf.writestr(f"{out_folder}/00_README.md", self._build_readme(index_entries))

        return zip_path

    async def _handle_attachments(
        self,
        node: Dict[str, Any],
        attachments_dir: str,
        zf: zipfile.ZipFile,
        downloader: Optional[Callable],
        current_dir: Optional[str] = None,
        caption_downloader: Optional[Callable] = None,
        progress_callback: Optional[Callable] = None,
        pct: float = 0,
    ) -> Tuple[Dict[str, str], Dict[str, str]]:
        mapping = {}
        doc_text_map = {}
        attachments = node.get("attachments", [])
        content_id = node.get("id")

        for att in attachments:
            att_id = att.get("id")
            filename = sanitize_filename(att.get("fileName", f"file_{att_id}"))
            zip_target_path = f"{attachments_dir}/{filename}"

            # Calculate relative path from document to attachment
            rel_path = f"./attachments/{filename}"

            original_url = att.get("originalUrl") or att.get("downloadUrl") or att_id
            mapping[original_url] = rel_path

            if downloader and content_id:
                download_key = att.get("downloadUrl") or att_id
                if download_key:
                    real_transcript_saved = False
                    lower_fn = filename.lower()
                    is_video_filename = any(lower_fn.endswith(ext) for ext in [".mp4", ".mov", ".m4v", ".webm", ".avi", ".mkv"])

                    data = None
                    try:
                        data = await downloader(self.course_id, content_id, download_key)

                        if is_video_filename and data and not _is_valid_video_bytes(data):
                            # download_attachment_bytes fell back to returning an
                            # HTML/error page (e.g. an unresolved LTI launch link) -
                            # don't save that mislabeled as a playable video file.
                            logger.warning(f"Attachment {filename} returned non-video data ({len(data)} bytes). Skipping.")
                            if progress_callback:
                                await progress_callback(f"Skipped {filename}: server returned an error page instead of the file", pct)
                            data = None
                    except FileTooLargeError as e:
                        size_str = _format_mb(e.size_bytes)
                        logger.warning(f"Skipped {filename}: exceeds size cap ({size_str})")
                        if progress_callback:
                            await progress_callback(f"Skipped {filename}: exceeds 300MB size limit ({size_str}) - extracting transcript only", pct)
                    except Exception as e:
                        logger.error(f"Failed to download attachment {filename}: {e}")
                        if progress_callback:
                            await progress_callback(f"Error downloading {filename}: download failed, skipping", pct)

                    if data:
                        written_bytes = _write_downloaded(zf, zip_target_path, data)
                        if written_bytes is not None:
                            extracted_text = extract_text_from_file_bytes(written_bytes, filename)
                            is_caption_data = (
                                written_bytes.startswith(b"WEBVTT")
                                or b"-->" in written_bytes[:500]
                                or lower_fn.endswith(".srt")
                                or lower_fn.endswith(".vtt")
                            )

                            if is_caption_data:
                                self._save_caption_and_transcript(node, written_bytes, lower_fn, zf, current_dir, doc_text_map)
                                real_transcript_saved = True
                            elif extracted_text and extracted_text.strip():
                                doc_text_map[filename] = extracted_text

                    # A Kaltura video's captions live as a separate asset, so fetch
                    # them independently rather than only as a fallback for when
                    # the video download fails, is too large, or returns non-video data.
                    if caption_downloader and not real_transcript_saved and _looks_like_kaltura_attachment(download_key, att):
                        try:
                            cap_data = await caption_downloader(self.course_id, content_id, download_key)
                            if cap_data:
                                self._save_caption_and_transcript(
                                    node, cap_data, "captions.srt", zf, current_dir, doc_text_map
                                )
                                if progress_callback:
                                    await progress_callback(f"Downloaded transcript for: {filename}", pct)
                        except Exception as e:
                            logger.warning(f"Failed to fetch Kaltura captions for {filename}: {e}")

        return mapping, doc_text_map

    def _save_caption_and_transcript(
        self,
        node: Dict[str, Any],
        data: bytes,
        lower_fn: str,
        zf: zipfile.ZipFile,
        current_dir: Optional[str],
        doc_text_map: Optional[Dict[str, str]] = None,
        video_basename: Optional[str] = None,
    ) -> None:
        """
        Saves a raw .srt/.vtt subtitle file plus its cleaned plain-text transcript.

        video_basename: when a real video file was also saved alongside this
        (e.g. "MyLecture" for "MyLecture.mp4"), also writes a copy of the raw
        subtitle named to match it exactly (e.g. "MyLecture.srt") - most media
        players (VLC, Windows Media Player, etc.) only auto-load a subtitle
        track when its filename matches the video's, not a suffixed variant.
        """
        clean_node_title = sanitize_filename(node.get("title", "Video"))
        sub_ext = ".vtt" if data.startswith(b"WEBVTT") or lower_fn.endswith(".vtt") else ".srt"
        target_dir = current_dir or self.root_folder_name

        # 1. Save raw subtitle file (.srt / .vtt)
        sub_filename = f"{clean_node_title}_subtitles{sub_ext}"
        zf.writestr(f"{target_dir}/{sub_filename}", data)

        # 1b. Also save a copy matching the video's filename so players auto-load it
        if video_basename:
            zf.writestr(f"{target_dir}/{video_basename}{sub_ext}", data)

        # 2. Save converted plain text transcript (.txt)
        raw_text = data.decode("utf-8", errors="ignore")
        clean_txt = clean_vtt_srt_transcript(raw_text)
        if clean_txt and clean_txt.strip():
            txt_filename = f"{clean_node_title}_transcript.txt"
            zf.writestr(f"{target_dir}/{txt_filename}", clean_txt)
            if doc_text_map is not None:
                doc_text_map[txt_filename] = clean_txt

    def _count_nodes(self, nodes: List[Dict[str, Any]]) -> int:
        count = 0
        for node in nodes:
            count += 1
            if node.get("children"):
                count += self._count_nodes(node["children"])
        return count

    async def build_raw_zip_package(
        self,
        content_tree: List[Dict[str, Any]],
        attachment_downloader: Optional[Callable[[str, str, str], Any]] = None,
        caption_downloader: Optional[Callable[[str, str, str], Any]] = None,
        embed_url_resolver: Optional[Callable[[str], Any]] = None,
        video_downloader: Optional[Callable[[str], Any]] = None,
        progress_callback: Optional[Callable[[str, float], Any]] = None,
    ) -> str:
        """
        Builds a ZIP package containing ONLY raw course files, PDFs, slides, and documents.
        Bypasses all Markdown conversion completely. Writes directly to a temp file on
        disk (see build_zip_package) and returns its path.
        """
        fd, zip_path = tempfile.mkstemp(suffix=".zip")
        os.close(fd)
        total_nodes = self._count_nodes(content_tree)
        safe_name = sanitize_filename(self.clean_course_name)
        root_dir = f"{safe_name}_RawFiles"

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            await self._process_raw_node_list(
                nodes=content_tree,
                current_dir=root_dir,
                zf=zf,
                attachment_downloader=attachment_downloader,
                caption_downloader=caption_downloader,
                embed_url_resolver=embed_url_resolver,
                video_downloader=video_downloader,
                progress_callback=progress_callback,
                processed_count=[0],
                total_nodes=total_nodes,
            )

        return zip_path

    async def build_raw_zip_package_streaming(
        self,
        topics,
        total_topics: int,
        attachment_downloader: Optional[Callable[[str, str, str], Any]] = None,
        caption_downloader: Optional[Callable[[str, str, str], Any]] = None,
        embed_url_resolver: Optional[Callable[[str], Any]] = None,
        video_downloader: Optional[Callable[[str], Any]] = None,
        progress_callback: Optional[Callable[[str, float], Any]] = None,
        course_id: Optional[str] = None,
        expand_fn: Optional[Callable] = None,
    ) -> str:
        """
        Same output as build_raw_zip_package, but `topics` is an async
        iterator that yields one top-level topic node at a time instead of
        a pre-built list of the whole course, and (when expand_fn is given)
        each topic's own folders are expanded and released one at a time
        too - see build_zip_package_streaming for why a single large topic
        needs this on top of topic-by-topic streaming, not just the latter.
        """
        fd, zip_path = tempfile.mkstemp(suffix=".zip")
        os.close(fd)
        safe_name = sanitize_filename(self.clean_course_name)
        root_dir = f"{safe_name}_RawFiles"

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            topic_num = 0
            async for top_node in topics:
                topic_num += 1
                # Rough, cheap size estimate - see build_zip_package_streaming's
                # identical use of _count_nodes for why this undercounts deep
                # nesting but is fine for interpolating progress.
                topic_node_count = self._count_nodes([top_node])
                await self._process_raw_node_list(
                    nodes=[top_node],
                    current_dir=root_dir,
                    zf=zf,
                    attachment_downloader=attachment_downloader,
                    caption_downloader=caption_downloader,
                    embed_url_resolver=embed_url_resolver,
                    video_downloader=video_downloader,
                    progress_callback=progress_callback,
                    processed_count=[0],
                    total_nodes=topic_node_count,
                    progress_base=((topic_num - 1) / max(1, total_topics)) * 100,
                    progress_span=(1 / max(1, total_topics)) * 100,
                    course_id=course_id,
                    expand_fn=expand_fn,
                )
                # top_node falls out of scope here, freeing its body/text
                # content before the next topic is fetched.

        return zip_path

    async def _process_raw_node_list(
        self,
        nodes: List[Dict[str, Any]],
        current_dir: str,
        zf: zipfile.ZipFile,
        attachment_downloader: Optional[Callable],
        progress_callback: Optional[Callable],
        processed_count: List[int],
        total_nodes: int,
        caption_downloader: Optional[Callable] = None,
        embed_url_resolver: Optional[Callable] = None,
        video_downloader: Optional[Callable] = None,
        progress_base: float = 0.0,
        progress_span: float = 100.0,
        course_id: Optional[str] = None,
        expand_fn: Optional[Callable] = None,
    ):
        # progress_base/progress_span let a caller processing one topic at a
        # time (build_raw_zip_package_streaming) blend this topic's 0-100%
        # progress into its slice of the whole course's progress bar,
        # instead of the 0-100% range always meaning "this whole course".
        # expand_fn (BlackboardClient.build_content_node), when given, turns
        # a raw child stub into a real node right before it's touched - see
        # _ensure_expanded - so a folder full of sub-folders is expanded and
        # released one at a time instead of all at once.
        for raw_node in nodes:
            node = await _ensure_expanded(raw_node, course_id, expand_fn)
            processed_count[0] += 1
            within_topic = min(processed_count[0] / max(1, total_nodes), 1.0)
            pct = progress_base + within_topic * progress_span
            if progress_callback:
                await progress_callback(f"Downloading raw files for: {node.get('title', 'Untitled')}", pct)

            title = sanitize_filename(node.get("title", "Untitled"))
            is_folder = node.get("isFolder", False)

            attachments = node.get("attachments", [])
            content_id = node.get("id")

            if is_folder:
                folder_path = f"{current_dir}/{title}"
                # Process attachments on the folder itself
                if attachments:
                    for att in attachments:
                        att_id = att.get("downloadUrl") or att.get("id")
                        filename = sanitize_filename(att.get("fileName", f"file_{att.get('id')}"))
                        zip_target_path = f"{folder_path}/{filename}"
                        if attachment_downloader and content_id and att_id:
                            if progress_callback:
                                await progress_callback(f"Downloading file: {filename}...", pct)
                            try:
                                data = await attachment_downloader(self.course_id, content_id, att_id)
                                if data:
                                    size_bytes = _content_len(data)
                                    _write_downloaded(zf, zip_target_path, data)
                                    if progress_callback:
                                        await progress_callback(f"Downloaded file: {filename} ({size_bytes} bytes)", pct)
                            except FileTooLargeError as e:
                                size_str = _format_mb(e.size_bytes)
                                logger.warning(f"Skipped {filename}: exceeds size cap ({size_str})")
                                if progress_callback:
                                    await progress_callback(f"Skipped {filename}: exceeds 300MB size limit ({size_str})", pct)
                            except Exception as e:
                                logger.error(f"Failed to download folder attachment {filename}: {e}")
                                if progress_callback:
                                    await progress_callback(f"Error downloading {filename}: download failed, skipping", pct)

                children = node.get("children", [])
                if children:
                    await self._process_raw_node_list(
                        nodes=children,
                        current_dir=folder_path,
                        zf=zf,
                        attachment_downloader=attachment_downloader,
                        progress_callback=progress_callback,
                        processed_count=processed_count,
                        total_nodes=total_nodes,
                        caption_downloader=caption_downloader,
                        embed_url_resolver=embed_url_resolver,
                        video_downloader=video_downloader,
                        progress_base=progress_base,
                        progress_span=progress_span,
                        course_id=course_id,
                        expand_fn=expand_fn,
                    )
            else:
                # Raw attachments placed directly into current_dir
                downloaded_any = False
                real_transcript_saved = False
                if attachments:
                    for att in attachments:
                        att_id = att.get("downloadUrl") or att.get("id")
                        filename = sanitize_filename(att.get("fileName", f"file_{att.get('id')}"))
                        zip_target_path = f"{current_dir}/{filename}"

                        if attachment_downloader and content_id and att_id:
                            if progress_callback:
                                await progress_callback(f"Downloading file: {filename}...", pct)

                            async def _try_caption_fallback(valid_video_flag):
                                nonlocal real_transcript_saved
                                if caption_downloader and not real_transcript_saved and _looks_like_kaltura_attachment(att_id, att):
                                    try:
                                        cap_data = await caption_downloader(self.course_id, content_id, att_id)
                                        if cap_data:
                                            video_basename = os.path.splitext(filename)[0] if valid_video_flag else None
                                            self._save_caption_and_transcript(node, cap_data, "captions.srt", zf, current_dir, video_basename=video_basename)
                                            real_transcript_saved = True
                                            if progress_callback:
                                                clean_node_title = sanitize_filename(node.get("title", "Video"))
                                                await progress_callback(f"Downloaded subtitles & transcript for: {clean_node_title}", pct)
                                    except Exception as e:
                                        logger.warning(f"Failed to fetch Kaltura captions for {filename}: {e}")

                            data = None
                            try:
                                data = await attachment_downloader(self.course_id, content_id, att_id)
                            except FileTooLargeError as e:
                                size_str = _format_mb(e.size_bytes)
                                logger.warning(f"Skipped {filename}: exceeds size cap ({size_str})")
                                if progress_callback:
                                    await progress_callback(f"Skipped {filename}: exceeds 300MB size limit ({size_str}) - extracting transcript only", pct)
                            except Exception as e:
                                logger.error(f"Failed to download raw attachment {filename}: {e}")
                                if progress_callback:
                                    await progress_callback(f"Error downloading {filename}: download failed, skipping", pct)

                            if data and _content_len(data) > 100:
                                lower_fn = filename.lower()
                                is_caption_data = isinstance(data, bytes) and (
                                    data.startswith(b"WEBVTT")
                                    or b"-->" in data[:500]
                                    or lower_fn.endswith(".srt")
                                    or lower_fn.endswith(".vtt")
                                )

                                if is_caption_data:
                                    # Save both raw .srt / .vtt subtitle file AND converted plain text .txt transcript!
                                    self._save_caption_and_transcript(node, data, lower_fn, zf, current_dir)
                                    downloaded_any = True
                                    real_transcript_saved = True
                                    if progress_callback:
                                        clean_node_title = sanitize_filename(node.get("title", "Video"))
                                        await progress_callback(f"Downloaded subtitles & transcript for: {clean_node_title}", pct)
                                else:
                                    is_video_file = any(lower_fn.endswith(ext) for ext in [".mp4", ".mov", ".m4v", ".webm", ".avi", ".mkv"])

                                    # Validate binary video headers to prevent saving HTML/XML 404 error pages as .mp4
                                    valid_video = (not is_video_file) or _is_valid_video_bytes(data)

                                    if valid_video:
                                        size_bytes = _content_len(data)
                                        _write_downloaded(zf, zip_target_path, data)
                                        downloaded_any = True
                                        if progress_callback:
                                            await progress_callback(f"Downloaded file: {filename} ({size_bytes} bytes)", pct)
                                    else:
                                        logger.warning(f"Attachment {filename} returned non-video/HTML error data ({_content_len(data)} bytes). Skipping saving invalid video file.")
                                        if progress_callback:
                                            await progress_callback(f"Skipped {filename}: server returned an error page instead of the file", pct)
                                        _discard(data)

                                    # The Kaltura caption/subtitle asset is separate from the video
                                    # asset itself, so fetch it independently rather than only as a
                                    # fallback for when the video download fails.
                                    await _try_caption_fallback(valid_video)
                            else:
                                # No data at all (too large, failed, or empty) - if this
                                # looks like a Kaltura video, still get at least the transcript.
                                await _try_caption_fallback(False)

                # Check if this item is a Kaltura / Panopto / Video item
                body = node.get("body") or node.get("description") or node.get("instructions") or ""
                handler = str(node.get("contentHandler", "")).lower()
                is_kaltura_or_video = (
                    "kaltura" in handler
                    or "video" in handler
                    or "panopto" in handler
                    or "kaltura" in body.lower()
                    or "panopto" in body.lower()
                    or any(att.get("isKaltura") for att in attachments)
                )

                if is_kaltura_or_video:
                    # Write interactive HTML video launcher, URL shortcut & README for Kaltura/video items
                    embed_url = ""
                    for att in attachments:
                        orig = att.get("originalUrl") or ""
                        dl = att.get("downloadUrl") or ""
                        for candidate_url in [orig, dl]:
                            if candidate_url and "playManifest" not in candidate_url and not candidate_url.startswith("undefined") and not candidate_url.startswith("javascript:"):
                                embed_url = candidate_url
                                break
                        if embed_url:
                            break

                    if not embed_url and body:
                        import re
                        for m in re.finditer(r'(?:href|src)=["\']([^"\']+)["\']', body):
                            cand = m.group(1) or ""
                            if cand and "playManifest" not in cand and not cand.startswith("undefined") and not cand.startswith("javascript:"):
                                embed_url = cand
                                break

                    if embed_url:
                        if embed_url.startswith("undefined") or embed_url.startswith("javascript:"):
                            embed_url = ""
                        elif embed_url.startswith("/"):
                            embed_url = f"{self.base_url}{embed_url}"

                    if not embed_url and content_id:
                        embed_url = f"{self.base_url}/webapps/blackboard/content/launchLink.jsp?course_id={self.course_id}&content_id={content_id}"

                    # embed_url is often an intermediate LTI launch link (e.g. .../execute/blti/launchLink)
                    # rather than the actual Kaltura URL - try to resolve the real one via redirects
                    # so the HTML launcher links straight to it and the entry id becomes visible.
                    if embed_url_resolver and embed_url:
                        try:
                            resolved = await embed_url_resolver(f"{embed_url} {body}")
                            if resolved:
                                embed_url = resolved
                        except Exception as e:
                            logger.debug(f"Could not resolve embed URL for {title}: {e}")

                    # Try downloading the actual MP4 bytes directly via Kaltura's own API
                    # using the resolved entry id - this needs no Blackboard session at all,
                    # only an entry id and a Kaltura instance that permits API-based serving.
                    # Done before captions so a matching-filename subtitle copy can be
                    # saved alongside it for player auto-load.
                    real_video_saved = False
                    if video_downloader:
                        search_text = f"{embed_url} {body}"
                        if _looks_like_kaltura_attachment(search_text, {}):
                            try:
                                video_data = await video_downloader(search_text)
                                if video_data and _content_len(video_data) > 10000:
                                    size_bytes = _content_len(video_data)
                                    _write_downloaded(zf, f"{current_dir}/{title}.mp4", video_data)
                                    real_video_saved = True
                                    if progress_callback:
                                        await progress_callback(f"Downloaded video: {title} ({size_bytes} bytes)", pct)
                                elif video_data:
                                    _discard(video_data)
                            except FileTooLargeError as e:
                                size_str = _format_mb(e.size_bytes)
                                logger.warning(f"Skipped video {title}: exceeds size cap ({size_str})")
                                if progress_callback:
                                    await progress_callback(f"Skipped video {title}: exceeds 300MB size limit ({size_str}) - extracting transcript only", pct)
                            except Exception as e:
                                logger.warning(f"Failed to download Kaltura video for {title}: {e}")
                                if progress_callback:
                                    await progress_callback(f"Error downloading video {title}: download failed, skipping", pct)

                    # Kaltura video items are often embedded via body HTML with no
                    # separate downloadable "attachment" entry, so the per-attachment
                    # caption fetch above never runs for them. Try again here using
                    # whatever Kaltura entry id can be found in the embed URL or body.
                    if caption_downloader and not real_transcript_saved:
                        search_text = f"{embed_url} {body}"
                        if _looks_like_kaltura_attachment(search_text, {}):
                            try:
                                cap_data = await caption_downloader(self.course_id, content_id, search_text)
                                if cap_data:
                                    video_basename = title if real_video_saved else None
                                    self._save_caption_and_transcript(node, cap_data, "captions.srt", zf, current_dir, video_basename=video_basename)
                                    real_transcript_saved = True
                                    if progress_callback:
                                        clean_node_title = sanitize_filename(node.get("title", "Video"))
                                        await progress_callback(f"Downloaded subtitles & transcript for: {clean_node_title}", pct)
                            except Exception as e:
                                logger.warning(f"Failed to fetch Kaltura captions for {title}: {e}")

                    if embed_url and not real_video_saved:
                        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title} - Video</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; margin: 0; padding: 20px; background: #0f172a; color: white; display: flex; flex-direction: column; align-items: center; justify-content: center; min-height: 90vh; }}
        h1 {{ margin-bottom: 20px; font-size: 1.5rem; text-align: center; }}
        .video-container {{ width: 100%; max-width: 960px; aspect-ratio: 16/9; background: #000; border-radius: 12px; overflow: hidden; box-shadow: 0 10px 25px rgba(0,0,0,0.5); }}
        iframe {{ width: 100%; height: 100%; border: none; }}
        .btn {{ margin-top: 20px; padding: 12px 24px; background: #2563eb; color: white; text-decoration: none; border-radius: 8px; font-weight: 600; display: inline-block; }}
        .btn:hover {{ background: #1d4ed8; }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    <div class="video-container">
        <iframe src="{embed_url}" allowfullscreen allow="autoplay; fullscreen; encrypted-media"></iframe>
    </div>
    <a class="btn" href="{embed_url}" target="_blank">Open Video in Browser</a>
</body>
</html>"""
                        zf.writestr(f"{current_dir}/{title}_Kaltura_Video.html", html_content)
                        zf.writestr(f"{current_dir}/{title}_Kaltura_Link.url", f"[InternetShortcut]\nURL={embed_url}\n")

                        readme_txt = f"""Kaltura Video: {node.get('title')}

Direct MP4 video file download is protected by NTU's Kaltura media server policies (requires active NTULearn login session).

How to watch this video:
1. Double-click '{title}_Kaltura_Video.html' to open and watch the embedded video in any web browser.
2. Or double-click '{title}_Kaltura_Link.url' to open the original video directly on NTULearn.
"""
                        zf.writestr(f"{current_dir}/{title}_README.txt", readme_txt)

                    # If a real .srt/.vtt caption was already fetched above, don't clobber
                    # it with the description text or the "no transcript" placeholder.
                    if not real_transcript_saved:
                        clean_node_title = sanitize_filename(node.get("title", "Video"))
                        txt_filename = f"{clean_node_title}_transcript.txt"
                        txt_target_path = f"{current_dir}/{txt_filename}"

                        clean_body_text = BeautifulSoup(body, "html.parser").get_text(separator="\n\n", strip=True) if body else ""
                        is_readme_text = any(phrase in clean_body_text for phrase in [
                            "Direct MP4 video file download is protected by",
                            "How to watch this video",
                            "Kaltura Video:",
                        ])

                        if clean_body_text and len(clean_body_text) > 10 and not is_readme_text:
                            zf.writestr(txt_target_path, clean_body_text)
                        else:
                            no_transcript_msg = f"""Title: {node.get('title')}
Course ID: {self.course_id}

Transcript Status: No .srt or .vtt subtitle file was uploaded or attached for this video by the instructor on NTULearn.

To watch the video with player controls in your browser:
1. Double-click '{title}_Kaltura_Video.html' to open and play the video.
2. Or double-click '{title}_Kaltura_Link.url' to view it directly on NTULearn.
"""
                            zf.writestr(txt_target_path, no_transcript_msg)

                elif not downloaded_any:
                    # Fallback for general content items with no attachments: save body HTML if present
                    if body and body.strip():
                        zf.writestr(f"{current_dir}/{title}.html", f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{title}</title></head><body><h1>{node.get('title')}</h1>{body}</body></html>")
